# app/services/pptx_service.py
from io import BytesIO
from fastapi import UploadFile
from pptx import Presentation

MAX_PPTX_FILE_SIZE = 5 * 1024 * 1024  # 5MB (기획 제한과 동일)

def extract_text_from_pptx(file: UploadFile) -> str:
    if file is None or not file.filename:
        raise ValueError("No file provided.")
    if not file.filename.lower().endswith(".pptx"):
        raise ValueError("The provided file is not a PPTX.")
    
    # size check
    file.file.seek(0, 2)
    size = file.file.tell()
    file.file.seek(0)
    if size > MAX_PPTX_FILE_SIZE:
        raise ValueError("The provided PPTX file is too large.")

    data = file.file.read()
    if not data:
        return ""
    
    prs = Presentation(BytesIO(data))

    lines: list[str] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        # 슬라이드 제목/본문 등 텍스트
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    lines.append(f"[Slide {slide_idx}] {t}")

    return "\n".join(lines).strip()
