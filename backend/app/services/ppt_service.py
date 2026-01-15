from __future__ import annotations

import logging
from io import BytesIO
from fastapi import UploadFile
from pptx import Presentation

logger = logging.getLogger(__name__)

MAX_PPT_FILE_SIZE = 10 * 1024 * 1024  # 10MB

def extract_text_from_ppt(file: UploadFile) -> str:
    """
    Extract text from a PPT/PPTX file.
    """
    if file is None or not file.filename:
        raise ValueError("No file provided.")
    
    if not file.filename.lower().endswith(".pptx"):
        raise ValueError("Only .pptx files are supported. Please save your .ppt file as .pptx.")
        
    file.file.seek(0, 2)
    file_size = file.file.tell()
    if file_size > MAX_PPT_FILE_SIZE:
        raise ValueError("The provided PPT file is too large.")
    file.file.seek(0)
    
    try:
        prs = Presentation(file.file)
        texts = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
                    
        merged_text = "\n\n".join(texts).strip()
        
        if not merged_text:
            raise ValueError("No text found in the provided PPT file.")
            
        return merged_text
        
    except Exception as e:
        logger.exception("PPT extract failed: %s", e)
        raise ValueError("Failed to extract text from PPT file.") from e
    finally:
        file.file.close()
