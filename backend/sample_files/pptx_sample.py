from pptx import Presentation
from pptx.util import Inches, Pt

def create_sample_pptx(filename: str = "sample.pptx"):
    prs = Presentation()

    # ---------- 슬라이드 1: 제목 슬라이드 ----------
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = "KNU MLA Sample PPTX"
    slide.placeholders[1].text = "PDF / PPTX 처리 테스트용 샘플"

    # ---------- 슬라이드 2: 본문 ----------
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = "강의 개요"
    content = slide.placeholders[1].text_frame
    content.text = "이 PPTX 파일은 API 테스트를 위한 샘플입니다."

    p = content.add_paragraph()
    p.text = "• PPTX → 텍스트 추출 테스트"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• 요약 및 번역 기능 확인"
    p.level = 1

    # ---------- 슬라이드 3: 자유 텍스트 ----------
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    textbox = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(8), Inches(2)
    )
    tf = textbox.text_frame
    tf.text = "자유 형식 텍스트 박스입니다."

    p = tf.add_paragraph()
    p.text = "여러 줄의 텍스트가 포함되어 있습니다."
    p.level = 1

    # ---------- 파일 저장 ----------
    prs.save(filename)
    print(f"✅ PPTX sample file created: {filename}")


if __name__ == "__main__":
    create_sample_pptx()
